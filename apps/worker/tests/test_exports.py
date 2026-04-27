"""Smoke tests for the real export builders.

Builds each artifact from the Kimpton Angler fixture and asserts the
resulting file exists with the expected structure (10 sheets / 8 slides /
non-empty PDF).
"""

from __future__ import annotations

import os
import shutil
import tempfile
import zipfile
from pathlib import Path

import pytest

os.environ.setdefault("DATABASE_URL", "sqlite+aiosqlite:///./fondok.db")


@pytest.fixture
def tmp_out() -> Path:
    d = Path(tempfile.mkdtemp(prefix="fondok-exports-"))
    yield d
    shutil.rmtree(d, ignore_errors=True)


def test_excel_builds(tmp_out: Path) -> None:
    """Excel builder produces a 10-sheet workbook with non-zero size."""
    from openpyxl import load_workbook

    from app.export import build_excel
    from app.export.fixtures import kimpton_model

    out = tmp_out / "model.xlsx"
    build_excel("kimpton-angler-2026", kimpton_model(), out)

    assert out.exists(), "xlsx not written"
    assert out.stat().st_size > 8_000, "xlsx is suspiciously small"

    wb = load_workbook(out, read_only=True)
    expected = {
        "Cover", "Assumptions", "Sources & Uses", "Operating Proforma",
        "Debt Schedule", "Returns", "Sensitivity", "Partnership",
        "Variance", "Market Comps",
    }
    assert set(wb.sheetnames) == expected, f"sheet mismatch: got {wb.sheetnames}"
    assert len(wb.sheetnames) == 10
    wb.close()


def test_memo_pdf_builds(tmp_out: Path) -> None:
    """WeasyPrint memo PDF is non-empty and has the PDF magic bytes."""
    pytest.importorskip(
        "weasyprint",
        reason="weasyprint requires system libs (cairo/pango)",
    )

    from app.export import build_memo_pdf
    from app.export.fixtures import kimpton_memo, kimpton_model

    out = tmp_out / "memo.pdf"
    build_memo_pdf(kimpton_memo(), kimpton_model(), out)

    assert out.exists(), "memo pdf not written"
    size = out.stat().st_size
    assert size > 4_000, f"pdf is suspiciously small ({size} bytes)"

    with out.open("rb") as fh:
        magic = fh.read(5)
    assert magic == b"%PDF-", f"file does not start with PDF magic, got {magic!r}"


def test_pptx_builds(tmp_out: Path) -> None:
    """Presentation has exactly 8 slides and is a valid pptx zip."""
    from pptx import Presentation

    from app.export import build_pptx
    from app.export.fixtures import kimpton_deal, kimpton_memo, kimpton_model

    out = tmp_out / "deck.pptx"
    build_pptx(kimpton_deal(), kimpton_model(), kimpton_memo(), out)

    assert out.exists(), "pptx not written"
    assert out.stat().st_size > 12_000, "pptx is suspiciously small"
    assert zipfile.is_zipfile(out), "pptx is not a valid zip archive"

    prs = Presentation(str(out))
    assert len(prs.slides) == 8, f"expected 8 slides, got {len(prs.slides)}"
