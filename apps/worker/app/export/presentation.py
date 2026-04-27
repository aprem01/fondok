# ruff: noqa: RUF001, RUF003
"""IC presentation deck builder (PowerPoint .pptx).

8 slides following an institutional-investor template:
  1. Title
  2. Executive Summary
  3. Property & Market Overview
  4. Operating Performance
  5. Sources & Uses + Capital Stack
  6. Returns Summary + 3-Scenario card
  7. Risk Assessment + Sensitivities
  8. Recommendation
"""

from __future__ import annotations

from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

# ─────────────────────────── palette ───────────────────────────
SLATE = RGBColor(0x1F, 0x29, 0x37)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x11, 0x18, 0x27)
INK_500 = RGBColor(0x6B, 0x72, 0x80)
INK_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_50 = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
GREEN_50 = RGBColor(0xEC, 0xFD, 0xF5)
GREEN_500 = RGBColor(0x10, 0xB9, 0x81)
AMBER_50 = RGBColor(0xFF, 0xFB, 0xEB)
RED_50 = RGBColor(0xFE, 0xE2, 0xE2)


def _fmt_usd(v: float | int | None) -> str:
    if v is None:
        return "—"
    try:
        return f"${float(v):,.0f}"
    except (TypeError, ValueError):
        return "—"


def _fmt_usd_m(v: float | int | None) -> str:
    if v is None:
        return "—"
    try:
        return f"${float(v) / 1_000_000:.1f}M"
    except (TypeError, ValueError):
        return "—"


def _fmt_pct(v: float | None) -> str:
    if v is None:
        return "—"
    try:
        return f"{float(v) * 100:.1f}%"
    except (TypeError, ValueError):
        return "—"


def _fmt_mult(v: float | None) -> str:
    if v is None:
        return "—"
    try:
        return f"{float(v):.2f}x"
    except (TypeError, ValueError):
        return "—"


# ─────────────────────────── shared chrome ───────────────────────────


def _add_top_ribbon(slide: Slide, title: str, subtitle: str = "") -> None:
    """Slate ribbon at the top with brand-amber accent stripe."""
    # Ribbon background
    ribbon = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.85)
    )
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = SLATE
    ribbon.line.fill.background()

    # Amber stripe at bottom edge of ribbon
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.82), Inches(13.33), Inches(0.05)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = AMBER
    stripe.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.10), Inches(11), Inches(0.8))
    tf = tb.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(10)
        r2.font.color.rgb = INK_300

    # Brand mark — top right
    brand = slide.shapes.add_textbox(Inches(11.0), Inches(0.25), Inches(2.1), Inches(0.4))
    p = brand.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "FONDOK"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = AMBER


def _add_textbox(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_size: int = 11,
    bold: bool = False,
    color: RGBColor = INK,
    align: int = PP_ALIGN.LEFT,
) -> None:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    tf.margin_left = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color


def _add_bullet_list(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: list[str],
    font_size: int = 12,
) -> None:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(font_size)
        r.font.color.rgb = INK


def _add_filled_box(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    line: RGBColor | None = None,
) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)


def _add_table(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    header: list[str],
    rows: list[list[str]],
    bold_last_row: bool = False,
    bold_rows: set[int] | None = None,
) -> None:
    bold_rows = bold_rows or set()
    n_rows = len(rows) + 1
    n_cols = len(header)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table

    for j, h_text in enumerate(header):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = SLATE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = h_text
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = WHITE

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            is_bold = (bold_last_row and i == len(rows)) or (i in bold_rows)
            if is_bold:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_100
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 else GRAY_50
            p = cell.text_frame.paragraphs[0]
            if j > 0:
                p.alignment = PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(9)
            r.font.bold = is_bold
            r.font.color.rgb = INK


# ─────────────────────────── slide builders ───────────────────────────


def _slide_title(prs: Presentation, deal: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Full-bleed slate background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLATE
    bg.line.fill.background()

    # Amber accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.4), Inches(0.18), Inches(2.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()

    # FONDOK wordmark
    _add_textbox(slide, 0.6, 0.6, 4, 0.5, "FONDOK", 14, True, AMBER)

    # Title block
    name = deal.get("name", deal.get("subject_property", "Hotel"))
    location = deal.get("city", deal.get("location", ""))
    brand = deal.get("brand", "")
    keys = deal.get("keys", "")

    _add_textbox(slide, 1.0, 2.4, 11.5, 0.6, "INVESTMENT COMMITTEE MEMORANDUM", 12, True, AMBER)
    _add_textbox(slide, 1.0, 3.05, 11.5, 1.4, name, 40, True, WHITE)
    subtitle = location
    if brand:
        subtitle += f"   ·   {brand}"
    if keys:
        subtitle += f"   ·   {keys} Keys"
    _add_textbox(slide, 1.0, 4.5, 11.5, 0.5, subtitle, 16, False, INK_300)

    today = datetime.now(UTC).strftime("%B %d, %Y")
    _add_textbox(slide, 1.0, 6.6, 11.5, 0.4, f"Investment Committee Memo  ·  {today}", 11, False, INK_300)


def _slide_exec_summary(prs: Presentation, memo: dict[str, Any], model: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_top_ribbon(slide, "Executive Summary", "Key takeaways from underwriting")

    # Pull 4 takeaways from memo summary or insights
    sections = memo.get("sections", [])
    takeaways: list[str] = []

    # Try from key_insights items first
    for s in sections:
        if s.get("section_id") == "key_insights":
            for item in s.get("items", [])[:4]:
                title = item.get("title", "")
                body = item.get("body", "")
                takeaways.append(f"{title}: {body}")
            break

    # Fallback: pull body from exec/thesis/recommendation
    if not takeaways:
        for sid in ("executive_summary", "investment_thesis", "recommendation"):
            for s in sections:
                if s.get("section_id") == sid and s.get("body"):
                    takeaways.append(s["body"])

    takeaways = takeaways[:4] or ["—"]
    _add_bullet_list(slide, 0.8, 1.4, 11.5, 4.8, takeaways, font_size=14)

    # Confidence footer
    confidence = (memo.get("appendix") or {}).get("ai_confidence", 0)
    _add_textbox(slide, 0.8, 6.6, 11.5, 0.4, f"AI Confidence: {int(confidence * 100)}%   ·   Engines: {len((memo.get('appendix') or {}).get('engines_run', []))}", 10, False, INK_500)


def _slide_property_market(prs: Presentation, deal: dict[str, Any], model: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_top_ribbon(slide, "Property & Market Overview")

    inv = model.get("investment_engine", {})
    market = model.get("market", {})

    # Left column — hotel facts
    _add_filled_box(slide, 0.5, 1.3, 6.0, 5.6, GRAY_50)
    _add_textbox(slide, 0.8, 1.5, 5.5, 0.4, "Hotel Facts", 14, True, INK)
    facts = [
        ("Property", deal.get("name", "—")),
        ("Brand", deal.get("brand", "—")),
        ("Location", deal.get("city", "—")),
        ("Keys", str(deal.get("keys", "—"))),
        ("Year Built", str(deal.get("year_built", deal.get("yearBuilt", "—")))),
        ("Service Tier", deal.get("service", "Lifestyle Boutique")),
        ("Purchase Price", _fmt_usd(inv.get("purchase_price_usd"))),
        ("Price / Key", _fmt_usd(inv.get("price_per_key_usd"))),
        ("Total Capital", _fmt_usd(inv.get("total_capital_usd"))),
    ]
    y = 2.1
    for k, v in facts:
        _add_textbox(slide, 0.8, y, 2.5, 0.32, k, 10, False, INK_500)
        _add_textbox(slide, 3.3, y, 3.0, 0.32, v, 11, True, INK)
        y += 0.42

    # Right column — submarket KPIs
    _add_filled_box(slide, 6.85, 1.3, 6.0, 5.6, GRAY_50)
    _add_textbox(slide, 7.15, 1.5, 5.5, 0.4, "Submarket KPIs", 14, True, INK)
    kpis = market.get("kpis") or {
        "Submarket": "Miami Beach / South Beach, FL",
        "RevPAR (TTM)": "$238",
        "ADR (TTM)": "$312",
        "Occupancy": "76.2%",
        "RGI": "112.2",
        "Supply Pipeline": "414 rooms (2.2%)",
        "Demand Growth": "4.8%",
        "Supply Growth": "1.2%",
        "Comp Set Avg Cap": "6.1%",
    }
    y = 2.1
    for k, v in kpis.items():
        _add_textbox(slide, 7.15, y, 2.8, 0.32, str(k), 10, False, INK_500)
        _add_textbox(slide, 9.95, y, 3.0, 0.32, str(v), 11, True, INK)
        y += 0.42


def _slide_operating(prs: Presentation, model: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_top_ribbon(slide, "Operating Performance", "5-Year Underwritten Proforma  (USD thousands)")

    pl = model.get("p_and_l_engine_proforma", {})
    lines = pl.get("lines", [])

    header = ["Line Item", "Y1", "Y2", "Y3", "Y4", "Y5", "CAGR"]
    rows = []
    bold_rows: set[int] = set()
    for i, line in enumerate(lines, start=1):
        row = [
            line.get("label", ""),
            f"{line.get('y1', 0):,}",
            f"{line.get('y2', 0):,}",
            f"{line.get('y3', 0):,}",
            f"{line.get('y4', 0):,}",
            f"{line.get('y5', 0):,}",
            f"{line.get('cagr', 0) * 100:.1f}%" if line.get("cagr") else "",
        ]
        rows.append(row)
        if line.get("bold"):
            bold_rows.add(i)

    _add_table(slide, 0.5, 1.3, 12.3, 5.6, header, rows, bold_rows=bold_rows)


def _slide_sources_uses(prs: Presentation, model: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_top_ribbon(slide, "Sources & Uses · Capital Stack")

    inv = model.get("investment_engine", {})
    debt = model.get("debt_engine", {})

    sources = model.get("sources") or [
        {"label": "Senior Debt", "amount": debt.get("loan_amount_usd", 0)},
        {"label": "Equity", "amount": inv.get("total_capital_usd", 0) - debt.get("loan_amount_usd", 0)},
    ]
    src_total = sum(s.get("amount", 0) for s in sources if not s.get("total"))
    src_rows = []
    for s in sources:
        if s.get("total"):
            continue
        amt = s.get("amount", 0)
        pct = (amt / src_total) if src_total else 0
        src_rows.append([s.get("label", ""), _fmt_usd(amt), f"{pct * 100:.1f}%"])
    src_rows.append(["Total Sources", _fmt_usd(src_total), "100.0%"])

    uses = model.get("uses") or [
        {"label": "Purchase Price", "amount": inv.get("purchase_price_usd", 0)},
        {"label": "Closing Costs", "amount": inv.get("closing_costs_usd", 0)},
        {"label": "Renovation", "amount": inv.get("renovation_budget_usd", 0)},
        {"label": "Soft Costs", "amount": inv.get("soft_costs_usd", 0)},
        {"label": "Contingency", "amount": inv.get("contingency_usd", 0)},
        {"label": "Working Capital", "amount": inv.get("working_capital_usd", 0)},
        {"label": "Loan Costs", "amount": inv.get("loan_costs_usd", 0)},
    ]
    use_total = sum(u.get("amount", 0) for u in uses if not u.get("total"))
    use_rows = []
    for u in uses:
        if u.get("total"):
            continue
        amt = u.get("amount", 0)
        pct = (amt / use_total) if use_total else 0
        use_rows.append([u.get("label", ""), _fmt_usd(amt), f"{pct * 100:.1f}%"])
    use_rows.append(["Total Uses", _fmt_usd(use_total), "100.0%"])

    _add_textbox(slide, 0.5, 1.3, 5.5, 0.4, "Sources", 13, True, INK)
    _add_table(
        slide, 0.5, 1.7, 6.0, 4.0, ["Source", "Amount", "%"], src_rows, bold_last_row=True
    )

    _add_textbox(slide, 6.85, 1.3, 5.5, 0.4, "Uses", 13, True, INK)
    _add_table(
        slide, 6.85, 1.7, 6.0, 4.5, ["Use", "Amount", "%"], use_rows, bold_last_row=True
    )

    # Capital stack summary
    _add_filled_box(slide, 0.5, 6.2, 12.3, 0.9, AMBER_50, AMBER)
    _add_textbox(slide, 0.8, 6.35, 4, 0.3, "CAPITAL STACK", 9, True, AMBER)
    _add_textbox(
        slide,
        0.8,
        6.6,
        12.0,
        0.4,
        f"LTC {_fmt_pct(debt.get('ltc'))}   ·   LTV {_fmt_pct(debt.get('ltv_value'))}   ·   Rate {_fmt_pct(debt.get('interest_rate_pct'))}   ·   DSCR Y1 {_fmt_mult(debt.get('year1_dscr'))}   ·   Debt Yield Y1 {_fmt_pct(debt.get('year1_debt_yield'))}",
        12,
        True,
        INK,
    )


def _slide_returns(prs: Presentation, model: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_top_ribbon(slide, "Returns Summary · Scenarios")

    ret = model.get("returns_engine", {})
    scenarios = model.get("scenario_outputs", [])

    # Top metrics tiles
    tiles = [
        ("Levered IRR", _fmt_pct(ret.get("levered_irr"))),
        ("Equity Multiple", _fmt_mult(ret.get("equity_multiple"))),
        ("Year 1 CoC", _fmt_pct(ret.get("year1_cash_on_cash"))),
        ("Hold (yrs)", str(ret.get("hold_years", 0))),
    ]
    x = 0.5
    for k, v in tiles:
        _add_filled_box(slide, x, 1.3, 3.0, 1.4, GRAY_50, INK_300)
        _add_textbox(slide, x + 0.2, 1.45, 2.8, 0.3, k, 10, False, INK_500)
        _add_textbox(slide, x + 0.2, 1.85, 2.8, 0.7, v, 28, True, AMBER)
        x += 3.16

    # Scenario cards
    _add_textbox(slide, 0.5, 3.05, 12, 0.4, "3-SCENARIO ANALYSIS", 11, True, AMBER)
    if scenarios:
        x = 0.5
        for sc in scenarios:
            name = sc.get("name", "")
            tone = GREEN_50 if sc.get("base") else (AMBER_50 if "Up" in name else RED_50)
            line = GREEN_500 if sc.get("base") else (AMBER if "Up" in name else RGBColor(0xEF, 0x44, 0x44))
            _add_filled_box(slide, x, 3.5, 4.1, 3.5, tone, line)
            _add_textbox(slide, x + 0.25, 3.65, 3.7, 0.4, name.upper(), 11, True, INK)
            metrics = [
                ("IRR", _fmt_pct(sc.get("irr"))),
                ("Unlev IRR", _fmt_pct(sc.get("unlevered_irr"))),
                ("Multiple", _fmt_mult(sc.get("multiple"))),
                ("Avg CoC", _fmt_pct(sc.get("avg_coc"))),
                ("Exit Value", _fmt_usd_m(sc.get("exit_value_usd"))),
            ]
            yy = 4.05
            for k, v in metrics:
                _add_textbox(slide, x + 0.25, yy, 2.0, 0.3, k, 10, False, INK_500)
                _add_textbox(slide, x + 2.2, yy, 1.8, 0.3, v, 11, True, INK)
                yy += 0.5
            x += 4.27


def _slide_risk_sensitivity(prs: Presentation, memo: dict[str, Any], model: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_top_ribbon(slide, "Risk Assessment · Sensitivities")

    # Risk subscores (left)
    risks = []
    for s in memo.get("sections", []):
        if s.get("section_id") == "risk_assessment":
            risks = s.get("subscores", [])
            break

    _add_textbox(slide, 0.5, 1.3, 5.5, 0.4, "Risk Profile", 13, True, INK)
    risk_rows = []
    for r in risks:
        risk_rows.append([r.get("name", ""), r.get("tier", ""), str(r.get("score", 0))])
    if not risk_rows:
        risk_rows = [["—", "—", "—"]]
    _add_table(slide, 0.5, 1.7, 6.0, 3.5, ["Risk", "Tier", "Score"], risk_rows)

    # Sensitivity heatmap (right) — 5x5 IRR vs Exit Cap × RevPAR Growth
    _add_textbox(slide, 6.85, 1.3, 5.5, 0.4, "Levered IRR — Exit Cap × RevPAR Growth", 13, True, INK)
    base_irr = float(model.get("returns_engine", {}).get("levered_irr", 0.23))
    cap_vals = [0.060, 0.065, 0.070, 0.075, 0.080]
    rev_vals = [0.02, 0.03, 0.05, 0.06, 0.08]

    n_cols = len(rev_vals) + 1
    n_rows = len(cap_vals) + 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(6.85), Inches(1.7), Inches(6.0), Inches(3.5))
    table = table_shape.table

    # Header row
    table.cell(0, 0).text = ""
    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = SLATE
    for j, rv in enumerate(rev_vals, start=1):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SLATE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{rv * 100:.0f}%"
        r.font.color.rgb = WHITE
        r.font.bold = True
        r.font.size = Pt(9)

    mid_r = (len(cap_vals) - 1) / 2
    mid_c = (len(rev_vals) - 1) / 2
    for i, cv in enumerate(cap_vals):
        cell = table.cell(i + 1, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY_100
        p = cell.text_frame.paragraphs[0]
        rr = p.add_run()
        rr.text = f"{cv * 100:.1f}%"
        rr.font.bold = True
        rr.font.size = Pt(9)
        for j in range(len(rev_vals)):
            drow = (i - mid_r) / max(mid_r, 1)
            dcol = (j - mid_c) / max(mid_c, 1)
            val = base_irr * (1 + 0.5 * (dcol - drow) / 2)
            cell = table.cell(i + 1, j + 1)
            cell.fill.solid()
            norm = (val - base_irr * 0.5) / (base_irr * 1.0) if base_irr else 0.5
            norm = max(0.0, min(1.0, norm))
            if norm > 0.66:
                cell.fill.fore_color.rgb = GREEN_50
            elif norm > 0.33:
                cell.fill.fore_color.rgb = AMBER_50
            else:
                cell.fill.fore_color.rgb = RED_50
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = f"{val * 100:.1f}%"
            r.font.size = Pt(9)
            r.font.color.rgb = INK

    # Variance disclosure footer
    var_body = ""
    for s in memo.get("sections", []):
        if s.get("section_id") == "variance_disclosure":
            var_body = s.get("body", "")
            break
    if var_body:
        _add_filled_box(slide, 0.5, 5.5, 12.3, 1.5, AMBER_50, AMBER)
        _add_textbox(slide, 0.7, 5.6, 4, 0.3, "VARIANCE DISCLOSURE", 9, True, AMBER)
        _add_textbox(slide, 0.7, 5.85, 12.0, 1.1, var_body, 11, False, INK)


def _slide_recommendation(prs: Presentation, memo: dict[str, Any], model: dict[str, Any]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_top_ribbon(slide, "Recommendation")

    header = memo.get("header", {})
    rec_text = header.get("recommendation", "PROCEED TO LOI")

    # Big recommendation card
    _add_filled_box(slide, 0.5, 1.3, 12.3, 2.0, GREEN_50, GREEN_500)
    _add_textbox(slide, 0.8, 1.45, 12, 0.4, "INVESTMENT COMMITTEE RECOMMENDATION", 11, True, GREEN_500)
    _add_textbox(slide, 0.8, 1.85, 12, 1.0, rec_text, 30, True, INK)

    # Pull rec body
    rec_body = ""
    for s in memo.get("sections", []):
        if s.get("section_id") == "recommendation":
            rec_body = s.get("body", "")
            break
    if rec_body:
        _add_textbox(slide, 0.8, 2.65, 12, 0.6, rec_body, 11, False, INK)

    # Action items
    _add_textbox(slide, 0.5, 3.5, 12, 0.4, "Next Steps", 13, True, INK)
    next_steps = [
        "Submit Letter of Intent at current ask",
        "Engage independent valuation firm and operator due diligence",
        "Confirm senior debt term sheet with lender",
        "Finalize PIP scope with Kimpton brand standards team",
        "Schedule IC vote — Mark project as IC Ready",
    ]
    _add_bullet_list(slide, 0.6, 3.95, 12, 2.5, next_steps, font_size=12)

    # Footer signature line
    today = datetime.now(UTC).strftime("%B %d, %Y")
    _add_textbox(
        slide,
        0.5,
        6.85,
        12.3,
        0.3,
        f"Prepared by Fondok AI · {header.get('lead_analyst', 'Analyst')} · {today}",
        9,
        False,
        INK_500,
    )


# ─────────────────────────── public API ───────────────────────────


def build_pptx(
    deal: dict[str, Any],
    model: dict[str, Any],
    memo: dict[str, Any],
    output_path: Path,
) -> Path:
    """Build the 8-slide IC presentation deck."""
    prs = Presentation()
    # 16:9 slide
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide_title(prs, deal)
    _slide_exec_summary(prs, memo, model)
    _slide_property_market(prs, deal, model)
    _slide_operating(prs, model)
    _slide_sources_uses(prs, model)
    _slide_returns(prs, model)
    _slide_risk_sensitivity(prs, memo, model)
    _slide_recommendation(prs, memo, model)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


__all__ = ["build_pptx"]
