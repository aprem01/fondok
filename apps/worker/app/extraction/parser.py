"""Document parsing — PDF (LlamaParse / PyMuPDF) and Excel (xlrd / openpyxl).

Routing logic
-------------
``parse_document`` dispatches on the filename extension:

* ``.pdf`` → PDF path. LlamaParse first when ``LLAMA_CLOUD_API_KEY`` is
  set; PyMuPDF + pdfplumber fallback otherwise.
* ``.xls`` (legacy BIFF8) → ``xlrd``. Each sheet becomes a
  ``ParsedPage`` so the LLM extractor sees one "page" per sheet.
* ``.xlsx`` (modern OOXML) → ``openpyxl``. Same per-sheet model.
* Anything else → ``ParseError``.

All paths emit a ``ParsedDocument`` so callers don't branch on backend.
``parser`` and ``metadata`` carry provenance so we can diagnose
extraction quality after the fact.

The legacy ``parse_pdf`` entry point still exists (callers haven't
migrated yet); it's now a thin alias that delegates to
``parse_document``.
"""

from __future__ import annotations

import asyncio
import hashlib
import logging
import os
import re
from datetime import UTC, datetime
from typing import Any

from .models import ParsedDocument, ParsedPage, ParseError

logger = logging.getLogger(__name__)


# ──────────────────────────── public ────────────────────────────


async def parse_document(file_bytes: bytes, filename: str) -> ParsedDocument:
    """Parse ``file_bytes`` into a ``ParsedDocument``.

    Dispatches on filename extension:
      * ``.pdf`` → LlamaParse (when configured) → PyMuPDF fallback.
      * ``.xls`` → ``xlrd`` (legacy BIFF8 — STR/CoStar exports use this).
      * ``.xlsx`` → ``openpyxl``.

    Runs the blocking parser in a thread to keep the FastAPI event loop
    responsive.
    """
    if not file_bytes:
        raise ParseError(f"empty file bytes for {filename}")

    content_hash = hashlib.sha256(file_bytes).hexdigest()
    ext = (filename or "").rsplit(".", 1)[-1].lower() if "." in (filename or "") else ""

    if ext == "pdf":
        api_key = os.environ.get("LLAMA_CLOUD_API_KEY")
        if api_key:
            try:
                return await asyncio.to_thread(
                    _parse_with_llamaparse,
                    file_bytes=file_bytes,
                    filename=filename,
                    content_hash=content_hash,
                    api_key=api_key,
                )
            except Exception as exc:  # noqa: BLE001 — fallback is the safety net
                logger.warning(
                    "parse_document: LlamaParse failed for %s (%s); falling back to PyMuPDF",
                    filename,
                    exc,
                )
        return await asyncio.to_thread(
            _parse_with_pymupdf,
            file_bytes=file_bytes,
            filename=filename,
            content_hash=content_hash,
        )

    if ext == "xls":
        return await asyncio.to_thread(
            _parse_with_xlrd,
            file_bytes=file_bytes,
            filename=filename,
            content_hash=content_hash,
        )

    # .xlsx and .xlsm both ride the openpyxl path — .xlsm is just a
    # macro-enabled OOXML workbook; we don't need the macros, only the
    # cell data. Rani's TTM uploads were arriving as .xlsm exports and
    # the old branch raised ParseError before extraction even started.
    if ext in {"xlsx", "xlsm"}:
        return await asyncio.to_thread(
            _parse_with_openpyxl,
            file_bytes=file_bytes,
            filename=filename,
            content_hash=content_hash,
        )

    # PowerPoint decks — OM teasers, market updates, and asset summaries
    # are routinely shared as .pptx (and occasionally .ppt). We pull the
    # text frames slide-by-slide via python-pptx; each slide becomes a
    # "page" so citations keep working. .ppt (legacy binary) is not
    # supported by python-pptx and surfaces as an actionable error.
    if ext == "pptx":
        return await asyncio.to_thread(
            _parse_with_pptx,
            file_bytes=file_bytes,
            filename=filename,
            content_hash=content_hash,
        )
    if ext == "ppt":
        raise ParseError(
            f"legacy binary .ppt ({filename}) is not supported — please "
            "re-save the deck as .pptx (File → Save As → PowerPoint "
            "Presentation) and re-upload."
        )

    raise ParseError(
        f"unsupported file extension '.{ext}' for {filename}; "
        "expected .pdf, .xls, .xlsx, .xlsm, or .pptx"
    )


async def parse_pdf(file_bytes: bytes, filename: str) -> ParsedDocument:
    """Backwards-compatible alias for ``parse_document``.

    Existing callers pass any document type through ``parse_pdf``;
    dispatch happens on filename so XLS/XLSX uploads now flow through
    the same single entry point.
    """
    return await parse_document(file_bytes, filename)


# ──────────────────────────── llamaparse ────────────────────────────


def _parse_with_llamaparse(
    *,
    file_bytes: bytes,
    filename: str,
    content_hash: str,
    api_key: str,
) -> ParsedDocument:
    """Synchronous LlamaParse path. Imported lazily.

    LlamaParse returns a list of Document objects (one per page in
    ``markdown`` mode). We harvest plain text from ``.text`` and pull
    pipe-tables out of the markdown for the ``tables`` channel.
    """
    try:
        from llama_parse import LlamaParse  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover — exercised when pkg missing
        raise ParseError(
            "llama-parse package not installed; cannot use LlamaParse"
        ) from exc

    import tempfile
    from pathlib import Path

    parser = LlamaParse(
        api_key=api_key,
        result_type="markdown",
        verbose=False,
    )

    # LlamaParse's load_data wants a path. Stash bytes in a tempfile
    # since we accept arbitrary uploads (including in-memory testing).
    with tempfile.NamedTemporaryFile(
        suffix=Path(filename).suffix or ".pdf", delete=False
    ) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        documents = parser.load_data(tmp_path)
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass

    pages: list[ParsedPage] = []
    for i, doc in enumerate(documents, start=1):
        text = getattr(doc, "text", "") or ""
        tables = _extract_markdown_tables(text)
        pages.append(
            ParsedPage(
                page_num=i,
                text=text,
                tables=tables,
                metadata={"source": "llamaparse"},
            )
        )

    return ParsedDocument(
        filename=filename,
        total_pages=len(pages),
        pages=pages,
        content_hash=content_hash,
        parsed_at=datetime.now(UTC),
        parser="llamaparse",
        metadata={"backend": "llamaparse"},
    )


def _extract_markdown_tables(markdown: str) -> list[list[list[str]]]:
    """Pull pipe-formatted tables out of a markdown blob.

    A markdown table looks like:
        | col1 | col2 |
        |------|------|
        | a    | b    |

    We collect contiguous lines starting with ``|`` and ending with
    ``|``. Header separators (``---``) are dropped. Cells are stripped.
    """
    tables: list[list[list[str]]] = []
    current: list[list[str]] = []

    def _flush() -> None:
        if current:
            tables.append([row for row in current if row])
            current.clear()

    for raw_line in markdown.splitlines():
        line = raw_line.strip()
        if line.startswith("|") and line.endswith("|"):
            cells = [c.strip() for c in line.strip("|").split("|")]
            # Skip pure separator rows.
            if all(re.fullmatch(r":?-+:?", c) for c in cells if c):
                continue
            current.append(cells)
        else:
            _flush()
    _flush()

    # Drop tables with fewer than 2 rows — they're almost always noise.
    return [t for t in tables if len(t) >= 2]


# ──────────────────────────── pymupdf fallback ────────────────────────────


def _parse_with_pymupdf(
    *,
    file_bytes: bytes,
    filename: str,
    content_hash: str,
) -> ParsedDocument:
    """Synchronous PyMuPDF + pdfplumber path.

    PyMuPDF (``fitz``) gives us per-page text fast. pdfplumber gives
    us table cells. Either may fail per-page; we log and continue —
    a partial parse is more useful than none for the underwriter.
    """
    try:
        import fitz  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise ParseError(
            "pymupdf is not installed; cannot parse PDFs"
        ) from exc

    pdfplumber_pdf: Any | None = None
    try:
        import pdfplumber  # type: ignore[import-untyped]

        import io

        pdfplumber_pdf = pdfplumber.open(io.BytesIO(file_bytes))
    except ImportError:
        logger.info(
            "pdfplumber not installed; skipping table extraction for %s",
            filename,
        )
        pdfplumber_pdf = None
    except Exception as exc:  # noqa: BLE001 — pdfplumber chokes on weird PDFs
        logger.warning(
            "pdfplumber failed to open %s: %s — proceeding without tables",
            filename,
            exc,
        )
        pdfplumber_pdf = None

    pages: list[ParsedPage] = []
    parser_label = "pymupdf+pdfplumber" if pdfplumber_pdf is not None else "pymupdf"

    try:
        with fitz.open(stream=file_bytes, filetype="pdf") as pdf:
            for i, page in enumerate(pdf, start=1):
                try:
                    text = page.get_text("text") or ""
                except Exception as exc:  # noqa: BLE001
                    logger.warning(
                        "pymupdf: failed to extract text on page %d of %s: %s",
                        i,
                        filename,
                        exc,
                    )
                    text = ""

                tables: list[list[list[str]]] = []
                if pdfplumber_pdf is not None and i <= len(pdfplumber_pdf.pages):
                    try:
                        plumber_page = pdfplumber_pdf.pages[i - 1]
                        raw_tables = plumber_page.extract_tables() or []
                        for tbl in raw_tables:
                            cleaned = [
                                [(cell or "").strip() for cell in row]
                                for row in tbl
                                if any((cell or "").strip() for cell in row)
                            ]
                            if cleaned:
                                tables.append(cleaned)
                    except Exception as exc:  # noqa: BLE001
                        logger.warning(
                            "pdfplumber: table extraction failed on page %d of %s: %s",
                            i,
                            filename,
                            exc,
                        )

                pages.append(
                    ParsedPage(
                        page_num=i,
                        text=text,
                        tables=tables,
                        metadata={"source": parser_label},
                    )
                )
    finally:
        if pdfplumber_pdf is not None:
            try:
                pdfplumber_pdf.close()
            except Exception:  # noqa: BLE001
                pass

    return ParsedDocument(
        filename=filename,
        total_pages=len(pages),
        pages=pages,
        content_hash=content_hash,
        parsed_at=datetime.now(UTC),
        parser=parser_label,
        metadata={"backend": parser_label},
    )


# ──────────────────────────── excel parsers ────────────────────────────
#
# Each sheet → one ``ParsedPage``. ``page_num`` reflects sheet order so
# downstream citations ("source_page=3") map to "the third sheet". The
# extractor's ``raw_text`` excerpt mechanism still works because each
# page carries the sheet's plain-text serialization plus the structured
# table grid.


def _format_xls_cell(value: Any) -> str:
    """Render an xlrd cell into a stable string.

    xlrd surfaces floats for numbers, ints for booleans, and ``""`` for
    empty cells. We normalize trailing-zero floats so "2.0" reads as "2"
    and dates as floats stay as floats (caller can parse).
    """
    if value is None or value == "":
        return ""
    if isinstance(value, float):
        return str(int(value)) if value.is_integer() else f"{value:g}"
    return str(value).strip()


def _xls_sheet_to_page(
    *,
    sheet_index: int,
    sheet_name: str,
    rows: list[list[str]],
) -> ParsedPage:
    """Serialize a 2D string grid into a ``ParsedPage``.

    ``text`` is a tab-separated dump (one row per line) so the LLM
    extractor sees the full grid as plain text. ``tables`` carries the
    same content as a structured grid so downstream callers that prefer
    cell-level access don't have to re-parse.
    """
    cleaned = [
        [c for c in row]
        for row in rows
        if any(c.strip() for c in row)
    ]
    text = "\n".join("\t".join(row) for row in cleaned)
    tables = [cleaned] if cleaned else []
    return ParsedPage(
        page_num=sheet_index,
        text=text,
        tables=tables,
        metadata={"source": "xls", "sheet_name": sheet_name},
    )


def _parse_with_xlrd(
    *,
    file_bytes: bytes,
    filename: str,
    content_hash: str,
) -> ParsedDocument:
    """Synchronous xlrd path for legacy ``.xls`` (BIFF8) files.

    STR / CoStar Trend Reports ship as 12-tab .xls workbooks; this is
    the only ingestion path that reaches them.
    """
    try:
        import xlrd  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise ParseError(
            "xlrd is not installed; cannot parse legacy .xls files"
        ) from exc

    try:
        wb = xlrd.open_workbook(file_contents=file_bytes)
    except Exception as exc:  # noqa: BLE001 — surface as ParseError
        raise ParseError(f"xlrd failed to open {filename}: {exc}") from exc

    pages: list[ParsedPage] = []
    for idx, sheet in enumerate(wb.sheets(), start=1):
        rows: list[list[str]] = []
        for r in range(sheet.nrows):
            row = [_format_xls_cell(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
            rows.append(row)
        pages.append(
            _xls_sheet_to_page(
                sheet_index=idx,
                sheet_name=sheet.name,
                rows=rows,
            )
        )

    return ParsedDocument(
        filename=filename,
        total_pages=len(pages),
        pages=pages,
        content_hash=content_hash,
        parsed_at=datetime.now(UTC),
        parser="xlrd",
        metadata={"backend": "xlrd", "sheet_count": len(pages)},
    )


def _parse_with_openpyxl(
    *,
    file_bytes: bytes,
    filename: str,
    content_hash: str,
) -> ParsedDocument:
    """Synchronous openpyxl path for modern ``.xlsx`` workbooks."""
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise ParseError(
            "openpyxl is not installed; cannot parse .xlsx files"
        ) from exc

    import io

    try:
        wb = load_workbook(
            io.BytesIO(file_bytes),
            data_only=True,
            read_only=True,
        )
    except Exception as exc:  # noqa: BLE001
        raise ParseError(f"openpyxl failed to open {filename}: {exc}") from exc

    pages: list[ParsedPage] = []
    try:
        for idx, sheet in enumerate(wb.worksheets, start=1):
            rows: list[list[str]] = []
            for row in sheet.iter_rows(values_only=True):
                rendered = [_format_xls_cell(v) for v in row]
                rows.append(rendered)
            pages.append(
                _xls_sheet_to_page(
                    sheet_index=idx,
                    sheet_name=sheet.title,
                    rows=rows,
                )
            )
    finally:
        try:
            wb.close()
        except Exception:  # noqa: BLE001
            pass

    return ParsedDocument(
        filename=filename,
        total_pages=len(pages),
        pages=pages,
        content_hash=content_hash,
        parsed_at=datetime.now(UTC),
        parser="openpyxl",
        metadata={"backend": "openpyxl", "sheet_count": len(pages)},
    )


def _parse_with_pptx(
    *,
    file_bytes: bytes,
    filename: str,
    content_hash: str,
) -> ParsedDocument:
    """Pull text frames + tables out of a .pptx deck via python-pptx.

    OM teasers, market updates, and asset summaries routinely ship as
    PowerPoint. We treat each slide as a "page" so the extractor's
    citation paths (slide N) keep working downstream. Tables on the
    slide get pulled as ``ParsedPage.tables`` for any structured
    line-item the extractor needs to ground against.

    Notes on the input:
    * Speaker notes are included (they often carry deal context the
      slide layout hides).
    * Grouped shapes are walked recursively — a "title + value" stat
      grouped inside a frame would otherwise be skipped.
    * Image-only slides surface as empty pages; the empty-envelope
      vs no-text differentiator downstream tells the user to OCR.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches  # noqa: F401 — keeps lint quiet
    except ImportError as exc:  # pragma: no cover
        raise ParseError(
            "python-pptx is not installed; cannot parse .pptx files"
        ) from exc

    import io

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:  # noqa: BLE001 — surface as ParseError
        raise ParseError(f"python-pptx failed to open {filename}: {exc}") from exc

    def _walk_shapes(shapes: Any) -> tuple[list[str], list[list[list[str]]]]:
        """Recursively collect text frames + tables from a shape tree.

        Returns ``(text_chunks, tables)``. PowerPoint allows shapes
        to be grouped, so we walk into ``shape_type == 6`` (GROUP)
        recursively rather than reading the top-level iterator only.
        """
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        for shape in shapes:
            # Grouped shapes — recurse.
            if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
                inner_t, inner_tb = _walk_shapes(shape.shapes)
                texts.extend(inner_t)
                tables.extend(inner_tb)
                continue
            # Text frames — paragraphs split into runs; join with spaces
            # so multi-run lines (e.g. bolded numbers + plain labels)
            # come through as readable strings.
            if getattr(shape, "has_text_frame", False):
                tf = shape.text_frame
                for para in tf.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
            # Tables — emit as a list-of-rows so the same renderer
            # the Excel parser uses can format them downstream.
            if getattr(shape, "has_table", False):
                tbl = shape.table
                rows: list[list[str]] = []
                for row in tbl.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    tables.append(rows)
        return texts, tables

    pages: list[ParsedPage] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts, tables = _walk_shapes(slide.shapes)
        # Speaker notes carry deal context the layout hides — include
        # them as a trailing block so the extractor can reference them.
        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                texts.append(f"[Speaker notes] {notes_text}")
        body = "\n".join(texts)
        pages.append(
            ParsedPage(
                page_num=slide_index,
                text=body,
                tables=tables,
                metadata={"source": "pptx", "slide_index": slide_index},
            )
        )

    return ParsedDocument(
        filename=filename,
        total_pages=len(pages),
        pages=pages,
        content_hash=content_hash,
        parsed_at=datetime.now(UTC),
        parser="python-pptx",
        metadata={"backend": "python-pptx", "slide_count": len(pages)},
    )


__all__ = ["parse_pdf", "parse_document"]
